#!/usr/bin/env python3
"""Build a lecture deck (.pptx) from canonical chapter Markdown.

Sibling of `build_chapter_html.mjs`: same input (`source/chapters/*.md`),
different output target. The Markdown remains authoritative — decks are
generated artifacts and are never hand-patched. Slide-specific wording,
extra slides, and slide-only images live in a per-chapter seed file
(`pipeline/slides/chNN.md`), which survives rebuilds.

Every slide carries a stable key plus a content fingerprint, written into
the speaker notes and into the deck model JSON. The key is what a later
ingest pass matches on; the fingerprint is the fallback when a key has been
deleted, and the signal that a slide's source prose has changed underneath
an override.

Output is written to `build/decks/` (gitignored). Decks are rebuilt, not
committed — a .pptx is a zip, so git stores a whole new copy on every save.

Usage:
    python pipeline/build_chapter_deck.py 6
    python pipeline/build_chapter_deck.py 6 --template pipeline/slides/psych101.potx
    python pipeline/build_chapter_deck.py 6 --render      # QA PNGs via LibreOffice
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
SOURCE_DIR = ROOT / "source" / "chapters"
SEED_DIR = ROOT / "pipeline" / "slides"
ASSET_DIR = ROOT / "slides"
BUILD_DIR = ROOT / "build" / "decks"
IMAGE_CACHE = BUILD_DIR / ".imgcache"

CHAPTER_SOURCES = {
    1: "ch01-history-approaches.md",
    2: "ch02-research-methods.md",
    3: "ch03-neuroscience-biological-bases.md",
    4: "ch04-sensation-perception.md",
    5: "ch05-consciousness.md",
    6: "ch06-sleep.md",
    7: "ch07-learning.md",
    8: "ch08-memory.md",
    9: "09-thinking-language-intelligence.md",
    10: "ch10-lifespan-development.md",
    11: "ch11-social-psychology.md",
    12: "ch12-emotion-stress-coping.md",
    13: "ch13-psychological-disorders-therapy.md",
}


# --------------------------------------------------------------------------
# CONFIG — the slide grammar. Change behaviour here, not in the code below.
# --------------------------------------------------------------------------

CONFIG = {
    # Sections dropped entirely. Matched against the `## ` heading text.
    "skip_sections": ["Further Reading", "References"],
    # Divider slides. "numbered-only" gives one to `## Section N:` headings and
    # lets the rest set slide titles instead — a divider reading "Learning
    # Objectives" in front of a slide titled "Learning Objectives" is noise.
    "section_dividers": "numbered-only",   # numbered-only | all | none
    # `###`/`####` subsections. "role-only" gives a slide to the labelled ones
    # (Do Not Confuse, Classic Study) and folds plain ones into slide titles.
    "subsection_slides": "role-only",      # role-only | all | none
    # Prose handling. "grouped" collects consecutive paragraphs in a subsection
    # onto one slide, one claim per bullet, with the full prose in the speaker
    # notes. "per-paragraph" gives each paragraph its own slide. "notes-only"
    # keeps prose off the slides entirely. The full paragraph always reaches the
    # notes regardless of mode.
    "prose_mode": "grouped",
    "prose_bullets_per_slide": 4,
    "prose_chars_per_bullet": 150,
    "prose_chars_on_slide": 300,
    # Learning objectives per slide.
    "objectives_per_slide": 3,
    # Key Terms: "grouped" (n per slide), "one-each", or "skip".
    "key_terms": "grouped",
    "key_terms_per_slide": 4,
    # Review questions: stem + options on the slide, answer in the notes.
    "review_questions": True,
    "answers_in_notes": True,
    # Callouts (> blockquotes) and #### Do Not Confuse / Classic Study blocks.
    "callout_slides": True,
    # "Try it yourself" lab pointers get their own slide.
    "lab_slides": True,
    # Markdown tables become real PowerPoint tables.
    "table_slides": True,
    # Data rows per slide, header repeated on each. PowerPoint auto-grows rows
    # to fit wrapped text, so a long table silently runs off the bottom rather
    # than shrinking — split it instead.
    "table_rows_per_slide": 4,
    # Figures.
    "figure_slides": True,
    "caption_on_slide": True,
    # Images: cap the long edge before embedding. 2560px covers a 4K projector.
    # "auto" re-encodes photographic/illustrative figures as JPEG and keeps flat
    # diagrams and anything with transparency as PNG — the chapter figures are
    # illustrations saved losslessly, which is where the weight comes from.
    "max_image_px": 2560,
    "image_format": "auto",
    "jpeg_quality": 88,
    # 16:9, matching docs/images/psych101_figure_style_guide.md.
    "slide_width_in": 13.333,
    "slide_height_in": 7.5,
}

# Layout indices in the default python-pptx template. A custom .potx should
# keep the same layout order, or these get remapped by name where possible.
LAYOUT = {
    "title": 0,
    "title_content": 1,
    "section": 2,
    "title_only": 5,
    "blank": 6,
}


# --------------------------------------------------------------------------
# Markdown parsing
# --------------------------------------------------------------------------

FIGURE_RE = re.compile(r"^!\[(?P<alt>.*?)\]\((?P<path>.*?)\)\s*$")
CAPTION_RE = re.compile(r"^\*(?P<caption>(?:Figure|Table)\s+[\d.]+\..*)\*\s*$")
HEADING_RE = re.compile(r"^(?P<hashes>#{1,6})\s+(?P<text>.*?)\s*$")
CALLOUT_RE = re.compile(r"^>\s*\*\*(?P<label>[^:*]+):?\*\*\s*(?P<rest>.*)$")
ORDERED_RE = re.compile(r"^(\d+)\.\s+(.*)$")
LAB_RE = re.compile(r"^\*\*Try it yourself:\*\*\s*(?P<rest>.*)$")


def slugify(text: str) -> str:
    text = re.sub(r"[^\w\s-]", "", text.lower())
    text = re.sub(r"[\s_]+", "-", text).strip("-")
    return re.sub(r"-{2,}", "-", text)


def section_slug(heading: str) -> str:
    """Short, typeable key stem — seeds are hand-edited, so
    `ch06-s1-p2` beats `ch06-section-1-circadian-rhythms-and-sleep-pressure-p2`."""
    numbered = re.match(r"Section\s+(\d+)", heading)
    if numbered:
        return f"s{numbered.group(1)}"
    return "-".join(slugify(heading).split("-")[:3])


def fingerprint(text: str) -> str:
    normalized = re.sub(r"\s+", " ", strip_inline(text)).strip().lower()
    return hashlib.sha1(normalized.encode("utf-8")).hexdigest()[:10]


def strip_inline(text: str) -> str:
    """Markdown inline markup -> plain text."""
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    return text.strip()


def inline_runs(text: str):
    """Split inline markup into (text, bold, italic) runs so bolded key terms
    stay bold on the slide instead of showing raw asterisks."""
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    runs = []
    token = re.compile(r"(\*\*.+?\*\*|(?<!\*)\*(?!\*).+?(?<!\*)\*(?!\*))", re.S)
    pos = 0
    for match in token.finditer(text):
        if match.start() > pos:
            runs.append((text[pos:match.start()], False, False))
        chunk = match.group(0)
        if chunk.startswith("**"):
            runs.append((chunk[2:-2], True, False))
        else:
            runs.append((chunk[1:-1], False, True))
        pos = match.end()
    if pos < len(text):
        runs.append((text[pos:], False, False))
    return [r for r in runs if r[0]]


def parse_blocks(lines):
    """Chapter Markdown -> a flat list of typed blocks."""
    blocks = []
    i = 0
    while i < len(lines):
        raw = lines[i]
        line = raw.rstrip("\n")
        stripped = line.strip()

        if not stripped or stripped == "---":
            i += 1
            continue

        if stripped.startswith(">") and stripped.lstrip("> ").startswith("*Drafting"):
            i += 1
            continue
        if stripped.startswith("> Drafting history"):
            i += 1
            continue

        heading = HEADING_RE.match(stripped)
        if heading:
            blocks.append({
                "type": "heading",
                "level": len(heading.group("hashes")),
                "text": heading.group("text"),
            })
            i += 1
            continue

        figure = FIGURE_RE.match(stripped)
        if figure:
            caption = ""
            if i + 1 < len(lines):
                cap = CAPTION_RE.match(lines[i + 1].strip())
                if cap:
                    caption = cap.group("caption")
                    i += 1
            blocks.append({
                "type": "figure",
                "alt": figure.group("alt"),
                "path": figure.group("path"),
                "caption": caption,
            })
            i += 1
            continue

        if stripped.startswith("|"):
            rows = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                if not all(re.fullmatch(r":?-{2,}:?", c) for c in cells if c):
                    rows.append(cells)
                i += 1
            if rows:
                blocks.append({"type": "table", "rows": rows})
            continue

        if stripped.startswith(">"):
            quote = []
            while i < len(lines) and lines[i].strip().startswith(">"):
                quote.append(lines[i].strip().lstrip(">").strip())
                i += 1
            joined = " ".join(q for q in quote if q)
            callout = CALLOUT_RE.match("> " + joined)
            if callout:
                blocks.append({
                    "type": "callout",
                    "label": callout.group("label").strip(),
                    "text": callout.group("rest").strip(),
                })
            else:
                blocks.append({"type": "callout", "label": "Note", "text": joined})
            continue

        if ORDERED_RE.match(stripped):
            items = []
            while i < len(lines) and ORDERED_RE.match(lines[i].strip()):
                items.append(ORDERED_RE.match(lines[i].strip()).group(2))
                i += 1
            blocks.append({"type": "ordered_list", "items": items})
            continue

        lab = LAB_RE.match(stripped)
        if lab:
            blocks.append({"type": "lab", "text": lab.group("rest")})
            i += 1
            continue

        para = [stripped]
        i += 1
        while i < len(lines):
            nxt = lines[i].strip()
            if (not nxt or nxt == "---" or nxt.startswith(("#", ">", "|", "!["))
                    or ORDERED_RE.match(nxt)):
                break
            para.append(nxt)
            i += 1
        blocks.append({"type": "paragraph", "text": " ".join(para)})

    return blocks


# --------------------------------------------------------------------------
# Slide model
# --------------------------------------------------------------------------

def first_sentence(text: str, minimum: int = 60) -> str:
    """Lead sentence of a paragraph, skipping abbreviation periods (H.M., et al.)
    by requiring a plausible minimum length."""
    plain = text.strip()
    pos = 0
    while True:
        match = re.search(r"[.!?](?:\s|$)", plain[pos:])
        if not match:
            return plain
        end = pos + match.end()
        if end >= minimum or end >= len(plain):
            return plain[:end].strip()
        pos = end


def truncate(text: str, limit: int) -> str:
    if len(text) <= limit:
        return text
    cut = text[:limit]
    stop = max(cut.rfind(". "), cut.rfind("; "), cut.rfind(" — "))
    if stop > limit * 0.5:
        return cut[:stop + 1].strip()
    return cut.rsplit(" ", 1)[0].strip() + "…"


class KeyMaker:
    """Stable-ish slide keys: section slug + kind + ordinal within section.

    Positional within a section, so inserting a paragraph shifts only the keys
    after it inside that one section — not the whole deck. The fingerprint
    recorded alongside each key is what lets an ingest pass re-match a slide
    whose key moved.
    """

    def __init__(self, chapter_num):
        self.chapter = f"ch{chapter_num:02d}"
        self.section = "front"
        self.counters = {}

    def set_section(self, slug):
        self.section = slug or "front"
        self.counters = {}

    def next(self, kind, explicit=None):
        if explicit:
            return f"{self.chapter}-{explicit}"
        n = self.counters.get(kind, 0) + 1
        self.counters[kind] = n
        return f"{self.chapter}-{self.section}-{kind}{n}"


def build_slide_model(chapter_num, blocks, config):
    slides = []
    keys = KeyMaker(chapter_num)
    section_title = None
    subsection = None
    mode = "front"
    pending_terms = []
    pending_question = None
    prose_buffer = []

    def add(kind, **fields):
        source_text = fields.pop("source_text", "")
        slide = {
            "kind": kind,
            "key": fields.pop("key", None) or keys.next(kind),
            "fingerprint": fingerprint(source_text or fields.get("title", "")),
            "section": section_title,
            "subsection": subsection,
        }
        slide.update(fields)
        slides.append(slide)
        return slide

    def flush_terms():
        nonlocal pending_terms
        if not pending_terms or config["key_terms"] == "skip":
            pending_terms = []
            return
        per = 1 if config["key_terms"] == "one-each" else config["key_terms_per_slide"]
        for start in range(0, len(pending_terms), per):
            group = pending_terms[start:start + per]
            add(
                "terms",
                title="Key Terms" if per > 1 else group[0][0],
                terms=group,
                key=keys.next("terms", explicit=f"term-{slugify(group[0][0])}"),
                source_text=" ".join(t[0] for t in group),
            )
        pending_terms = []

    def flush_question():
        nonlocal pending_question
        if pending_question and pending_question.get("stem"):
            add("question", **pending_question)
        pending_question = None

    def flush_prose():
        """Consecutive paragraphs become one slide with one claim per bullet;
        the full prose always lands in the notes."""
        nonlocal prose_buffer
        buffered, prose_buffer = prose_buffer, []
        if not buffered:
            return
        style = config["prose_mode"]
        full = "\n\n".join(strip_inline(t) for t in buffered)

        if style == "notes-only":
            if slides:
                previous = slides[-1]
                previous["notes"] = (
                    (previous.get("notes", "") + "\n\n" + full).strip()
                )
            return

        per = 1 if style == "per-paragraph" else config["prose_bullets_per_slide"]
        for start in range(0, len(buffered), per):
            group = buffered[start:start + per]
            if per == 1:
                bullets = [truncate(group[0], config["prose_chars_on_slide"])]
            else:
                bullets = [
                    truncate(first_sentence(t), config["prose_chars_per_bullet"])
                    for t in group
                ]
            add("prose",
                title=subsection or section_title or "",
                bullets=bullets,
                notes="\n\n".join(strip_inline(t) for t in group),
                key=keys.next("p"),
                source_text=" ".join(group))

    for block in blocks:
        btype = block["type"]

        if btype == "heading":
            flush_prose()
            level, text = block["level"], block["text"]

            if level == 1:
                add("title", title=text, key=f"ch{chapter_num:02d}-title",
                    source_text=text)
                continue

            if level == 2:
                flush_terms()
                flush_question()
                if any(text.startswith(s) for s in config["skip_sections"]):
                    mode = "skip"
                    continue
                section_title = text
                subsection = None
                keys.set_section(section_slug(text))

                if text == "Learning Objectives":
                    mode = "objectives"
                elif text == "Key Terms":
                    mode = "terms"
                elif text == "Review Questions":
                    mode = "questions"
                elif text == "Misconception Opener":
                    mode = "opener"
                else:
                    mode = "body"

                style = config["section_dividers"]
                if style == "all" or (style == "numbered-only"
                                      and text.startswith("Section")):
                    add("section", title=text, key=keys.next("divider"),
                        source_text=text)
                continue

            if level in (3, 4) and mode not in ("skip", "terms", "questions"):
                subsection = text
                role = None
                for marker in ("Do Not Confuse", "Classic Study"):
                    if text.startswith(marker):
                        role = marker
                style = config["subsection_slides"]
                if style == "all" or (style == "role-only" and role):
                    add("subsection", title=text, role=role,
                        key=keys.next("sub"), source_text=text)
                continue
            continue

        if mode == "skip":
            continue

        if mode == "objectives" and btype == "ordered_list":
            per = config["objectives_per_slide"]
            items = [strip_inline(x) for x in block["items"]]
            for start in range(0, len(items), per):
                group = items[start:start + per]
                add("objectives", title="Learning Objectives",
                    items=group, numbered_from=start + 1,
                    key=keys.next("obj"), source_text=" ".join(group))
            continue

        if mode == "terms" and btype == "paragraph":
            match = re.match(r"^\*\*(?P<term>.+?)\*\*\s*[—-]\s*(?P<definition>.*)$",
                             block["text"])
            if match:
                pending_terms.append(
                    (match.group("term"), strip_inline(match.group("definition")))
                )
            continue

        if mode == "questions":
            if btype == "paragraph":
                text = block["text"]
                q = re.match(r"^\*\*(?P<num>\d+)\.\*\*\s*(?P<stem>.*)$", text)
                if q:
                    flush_question()
                    pending_question = {
                        "title": f"Review Question {q.group('num')}",
                        "stem": strip_inline(q.group("stem")),
                        "options": [],
                        "answer": "",
                        "key": keys.next("rq", explicit=f"rq-{q.group('num')}"),
                        "source_text": q.group("stem"),
                    }
                    continue
                if pending_question is not None:
                    if text.startswith("*Answer:"):
                        pending_question["answer"] = strip_inline(text)
                        flush_question()
                        continue
                    for opt in re.findall(r"([a-d])\)\s*([^\n]+?)(?=\s+[a-d]\)|$)", text):
                        pending_question["options"].append(
                            f"{opt[0]}) {strip_inline(opt[1])}"
                        )
            continue

        if btype == "figure" and config["figure_slides"]:
            flush_prose()
            fig_no = ""
            fig_match = re.match(r"Figure\s+([\d.]+)\s*:?\s*(?P<rest>.*)",
                                 block["alt"])
            lead = ""
            if fig_match:
                fig_no = fig_match.group(1)
                lead = re.split(r"(?<=[a-z])[.:;]\s", fig_match.group("rest"))[0]
            # The claim makes a better slide title than the figure number, which
            # is only useful for pointing back at the book.
            title = truncate(lead, 80) if lead else (
                f"Figure {fig_no}" if fig_no else "Figure")
            add("figure", title=title, image=block["path"], alt=block["alt"],
                figure_number=fig_no,
                caption=block["caption"],
                key=keys.next("fig", explicit=f"fig-{fig_no.replace('.', '-')}"
                              if fig_no else None),
                source_text=block["alt"])
            continue

        if btype == "table" and config["table_slides"]:
            flush_prose()
            header, body_rows = block["rows"][0], block["rows"][1:]
            per = max(1, config["table_rows_per_slide"])
            chunks = [body_rows[n:n + per] for n in range(0, len(body_rows), per)] \
                or [[]]
            base_title = subsection or section_title or ""
            for part, chunk in enumerate(chunks, start=1):
                title = base_title if len(chunks) == 1 else \
                    f"{base_title} ({part}/{len(chunks)})"
                add("table", title=title, rows=[header] + chunk,
                    key=keys.next("table"),
                    source_text=" ".join(header) + " ".join(c[0] for c in chunk))
            continue

        if btype == "callout" and config["callout_slides"]:
            flush_prose()
            add("callout", title=block["label"], body=strip_inline(block["text"]),
                key=keys.next("prompt"), source_text=block["text"])
            continue

        if btype == "lab" and config["lab_slides"]:
            flush_prose()
            link = re.search(r"\[(?P<name>[^\]]+)\]\((?P<url>[^)]+)\)", block["text"])
            add("lab", title=link.group("name") if link else "Try it yourself",
                url=link.group("url") if link else "",
                body=strip_inline(block["text"]),
                key=keys.next("lab"), source_text=block["text"])
            continue

        if btype == "paragraph":
            text = block["text"]
            if text.startswith("*") and text.endswith("*") and mode == "opener":
                flush_prose()
                add("quote", title="", body=strip_inline(text),
                    key=keys.next("quote"), source_text=text)
                continue
            plain = strip_inline(text)
            # Short fragments and lead-ins ("By the end of this chapter, you
            # should be able to:") introduce the next block rather than carrying
            # a claim of their own.
            if len(plain) < 40 or plain.endswith(":"):
                continue
            prose_buffer.append(text)
            if (config["prose_mode"] == "grouped"
                    and len(prose_buffer) >= config["prose_bullets_per_slide"]):
                flush_prose()
            continue

    flush_prose()
    flush_terms()
    flush_question()
    return slides


# --------------------------------------------------------------------------
# Seed overrides
# --------------------------------------------------------------------------

def load_seed(path: Path):
    """Per-chapter overrides. Format:

        ## <slide-key>
        title: Replacement title
        skip: true
        image: assets/rat-in-skinner-box.png
        insert-after: ch06-fig-6-2
        body:
          Multi-line body text, indented.
    """
    if not path.exists():
        return {}
    overrides = {}
    current_key = None
    current_field = None
    for raw in path.read_text(encoding="utf-8").splitlines():
        heading = re.match(r"^##\s+(\S+)\s*$", raw)
        if heading:
            current_key = heading.group(1)
            overrides[current_key] = {}
            current_field = None
            continue
        if current_key is None:
            continue
        field = re.match(r"^(?P<name>[a-z-]+):\s*(?P<value>.*)$", raw)
        if field and not raw.startswith((" ", "\t")):
            name, value = field.group("name"), field.group("value").strip()
            if value:
                overrides[current_key][name] = (
                    True if value == "true" else False if value == "false" else value
                )
                current_field = None
            else:
                current_field = name
                overrides[current_key][name] = ""
            continue
        if current_field and raw.strip():
            existing = overrides[current_key][current_field]
            overrides[current_key][current_field] = (
                existing + ("\n" if existing else "") + raw.strip()
            )
    return overrides


def apply_seed(slides, overrides):
    """Overrides win over generated content. Unmatched keys are reported, not
    silently dropped — an orphaned override usually means a chapter edit moved
    the slide it was keyed to."""
    by_key = {s["key"]: s for s in slides}
    applied, orphaned, inserted = [], [], []

    for key, fields in overrides.items():
        if key in by_key:
            slide = by_key[key]
            if fields.get("skip") is True:
                slide["_skip"] = True
            for name, value in fields.items():
                if name in ("skip", "insert-after"):
                    continue
                slide[name] = value
                slide["_overridden"] = True
            applied.append(key)
        elif fields.get("insert-after"):
            anchor = fields["insert-after"]
            if anchor in by_key:
                new = {
                    "kind": fields.get("kind", "prose"),
                    "key": key,
                    "fingerprint": fingerprint(fields.get("body", key)),
                    "section": by_key[anchor].get("section"),
                    "subsection": by_key[anchor].get("subsection"),
                    "_seeded": True,
                }
                new.update({k: v for k, v in fields.items() if k != "insert-after"})
                slides.insert(slides.index(by_key[anchor]) + 1, new)
                inserted.append(key)
            else:
                orphaned.append(key)
        else:
            orphaned.append(key)

    return [s for s in slides if not s.get("_skip")], applied, inserted, orphaned


# --------------------------------------------------------------------------
# Images
# --------------------------------------------------------------------------

def resolve_image(path_str: str, chapter_num: int):
    """Chapter Markdown writes figure paths for the generated HTML's location,
    not its own. Try the plausible roots, then fall back to filename lookup."""
    name = Path(path_str).name
    tail = re.sub(r"^(\.\./)+", "", path_str)
    candidates = [
        ROOT / "docs" / tail,
        ROOT / tail,
        ROOT / "docs" / "images" / f"ch{chapter_num:02d}" / name,
        ASSET_DIR / f"ch{chapter_num:02d}" / tail,
        ASSET_DIR / f"ch{chapter_num:02d}" / "assets" / name,
    ]
    for candidate in candidates:
        if candidate.exists():
            return candidate
    matches = list((ROOT / "docs" / "images").rglob(name))
    return matches[0] if matches else None


def choose_format(img, config):
    """Flat diagrams compress well losslessly; illustrations and photographs do
    not, and that is where deck weight comes from. Colour count separates them
    cheaply. Transparency always forces PNG."""
    if config["image_format"] in ("png", "jpeg"):
        return config["image_format"]
    if img.mode in ("RGBA", "LA", "P") and "transparency" in img.info:
        return "png"
    colors = img.convert("RGB").getcolors(maxcolors=20000)
    return "png" if colors is not None else "jpeg"


def prepare_image(src: Path, config):
    """Cap the long edge and re-encode into the build cache. The originals in
    docs/images are untouched; only the embedded copy is resized."""
    IMAGE_CACHE.mkdir(parents=True, exist_ok=True)
    cap = config["max_image_px"]
    stamp = hashlib.sha1(
        f"{src}|{src.stat().st_mtime_ns}|{cap}|{config['image_format']}"
        f"|{config['jpeg_quality']}".encode("utf-8")
    ).hexdigest()[:12]
    for suffix in (".jpg", ".png"):
        cached = IMAGE_CACHE / f"{src.stem}-{stamp}{suffix}"
        if cached.exists():
            return cached, src.stat().st_size, cached.stat().st_size

    with Image.open(src) as opened:
        img = opened.convert("RGBA") if opened.mode in ("P", "LA") else opened.copy()
        if max(img.size) > cap:
            ratio = cap / max(img.size)
            img = img.resize(
                (max(1, int(img.width * ratio)), max(1, int(img.height * ratio))),
                Image.LANCZOS,
            )
        fmt = choose_format(opened, config)

    if fmt == "jpeg":
        if img.mode != "RGB":
            background = Image.new("RGB", img.size, (255, 255, 255))
            background.paste(img, mask=img.split()[-1] if img.mode == "RGBA" else None)
            img = background
        cached = IMAGE_CACHE / f"{src.stem}-{stamp}.jpg"
        img.save(cached, "JPEG", quality=config["jpeg_quality"], optimize=True,
                 progressive=True)
    else:
        cached = IMAGE_CACHE / f"{src.stem}-{stamp}.png"
        img.save(cached, "PNG", optimize=True)
    return cached, src.stat().st_size, cached.stat().st_size


# --------------------------------------------------------------------------
# PPTX rendering
# --------------------------------------------------------------------------

def layout_for(prs, name):
    idx = LAYOUT[name]
    if idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def suppress_bullet(para):
    """Placeholders auto-bullet. When the text carries its own numbering, the
    two collide as '• 4. Compare…' — drop the glyph."""
    pPr = para._p.get_or_add_pPr()
    for existing in pPr.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buChar"
    ):
        pPr.remove(existing)
    pPr.append(pPr.makeelement(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}buNone", {}
    ))


def set_text(frame, blocks, size=20, bold=False, align=PP_ALIGN.LEFT,
             bullets=True):
    frame.word_wrap = True
    frame.clear()
    for n, block in enumerate(blocks):
        para = frame.paragraphs[0] if n == 0 else frame.add_paragraph()
        para.alignment = align
        para.space_after = Pt(10)
        if not bullets:
            suppress_bullet(para)
        for text, is_bold, is_italic in inline_runs(block):
            run = para.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold or is_bold
            run.font.italic = is_italic


def add_notes(slide, body, key, fp):
    frame = slide.notes_slide.notes_text_frame
    lines = []
    if body:
        lines.append(body.strip())
        lines.append("")
    lines.append(f"[slide-key: {key}]")
    lines.append(f"[fingerprint: {fp}]")
    frame.text = "\n".join(lines)


def render_deck(chapter_num, chapter_title, slides, config, template=None):
    prs = Presentation(str(template)) if template else Presentation()
    prs.slide_width = Inches(config["slide_width_in"])
    prs.slide_height = Inches(config["slide_height_in"])
    width, height = prs.slide_width, prs.slide_height
    margin = Inches(0.7)
    content_width = width - margin * 2

    stats = {"images": 0, "bytes_before": 0, "bytes_after": 0, "missing": []}

    for slide_def in slides:
        kind = slide_def["kind"]
        notes = slide_def.get("notes", "")

        if kind == "title":
            slide = prs.slides.add_slide(layout_for(prs, "title"))
            slide.shapes.title.text = slide_def.get("title", chapter_title)
            if len(slide.placeholders) > 1:
                try:
                    slide.placeholders[1].text = "Psych 101"
                except (KeyError, IndexError):
                    pass

        elif kind == "section":
            slide = prs.slides.add_slide(layout_for(prs, "section"))
            slide.shapes.title.text = slide_def.get("title", "")

        elif kind == "subsection":
            slide = prs.slides.add_slide(layout_for(prs, "section"))
            role = slide_def.get("role")
            title = slide_def.get("title", "")
            slide.shapes.title.text = title
            if role:
                notes = f"[{role}] " + (notes or "")

        elif kind == "quote":
            slide = prs.slides.add_slide(layout_for(prs, "blank"))
            box = slide.shapes.add_textbox(
                margin, Inches(2.2), content_width, Inches(3.0)
            )
            set_text(box.text_frame, [slide_def.get("body", "")],
                     size=32, align=PP_ALIGN.CENTER)
            for para in box.text_frame.paragraphs:
                for run in para.runs:
                    run.font.italic = True

        elif kind == "objectives":
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "Learning Objectives")
            start = slide_def.get("numbered_from", 1)
            body = slide.placeholders[1].text_frame
            set_text(
                body,
                [f"{start + n}. {item}" for n, item in enumerate(slide_def["items"])],
                size=20, bullets=False,
            )

        elif kind == "terms":
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "Key Terms")
            body = slide.placeholders[1].text_frame
            set_text(
                body,
                [f"**{term}** — {definition}"
                 for term, definition in slide_def["terms"]],
                size=16,
            )

        elif kind == "question":
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "Review Question")
            body = slide.placeholders[1].text_frame
            set_text(body, [slide_def["stem"]] + slide_def.get("options", []),
                     size=18)
            if config["answers_in_notes"]:
                notes = (slide_def.get("answer", "") + "\n\n" + (notes or "")).strip()

        elif kind == "table":
            slide = prs.slides.add_slide(layout_for(prs, "title_only"))
            slide.shapes.title.text = slide_def.get("title", "")
            rows = slide_def["rows"]
            n_rows, n_cols = len(rows), max(len(r) for r in rows)
            top = Inches(1.9)
            table_height = min(height - top - margin, Inches(0.5) * n_rows)
            shape = slide.shapes.add_table(
                n_rows, n_cols, margin, top, content_width, table_height
            )
            longest = max(
                (len(strip_inline(c)) for row in rows for c in row), default=0
            )
            cell_size = 13 if longest < 90 else 11
            for r, row in enumerate(rows):
                for c in range(n_cols):
                    cell = shape.table.cell(r, c)
                    cell.text = strip_inline(row[c]) if c < len(row) else ""
                    cell.margin_top = Pt(3)
                    cell.margin_bottom = Pt(3)
                    for para in cell.text_frame.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(cell_size)
                            run.font.bold = r == 0

        elif kind == "figure":
            slide = prs.slides.add_slide(layout_for(prs, "title_only"))
            slide.shapes.title.text = slide_def.get("title", "")
            src = resolve_image(slide_def["image"], chapter_num)
            if src is None:
                stats["missing"].append(slide_def["image"])
                box = slide.shapes.add_textbox(margin, Inches(3), content_width,
                                               Inches(1))
                set_text(box.text_frame,
                         [f"[missing image: {slide_def['image']}]"], size=16)
            else:
                prepared, before, after = prepare_image(src, config)
                stats["images"] += 1
                stats["bytes_before"] += before
                stats["bytes_after"] += after
                caption_space = Inches(1.0) if config["caption_on_slide"] else Inches(0.3)
                top = Inches(1.6)
                avail_h = height - top - caption_space - Inches(0.3)
                with Image.open(prepared) as img:
                    iw, ih = img.size
                scale = min(content_width / iw, avail_h / ih)
                draw_w, draw_h = Emu(int(iw * scale)), Emu(int(ih * scale))
                slide.shapes.add_picture(
                    str(prepared), Emu(int((width - draw_w) / 2)), top,
                    width=draw_w, height=draw_h,
                )
                if config["caption_on_slide"] and slide_def.get("caption"):
                    cap_box = slide.shapes.add_textbox(
                        margin, height - caption_space, content_width, Inches(0.8)
                    )
                    set_text(cap_box.text_frame, [slide_def["caption"]], size=11)
                notes = (slide_def.get("alt", "") + "\n\n" + (notes or "")).strip()

        elif kind == "callout":
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "")
            set_text(slide.placeholders[1].text_frame,
                     [slide_def.get("body", "")], size=22)

        elif kind == "lab":
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "Try it yourself")
            lines = [slide_def.get("body", "")]
            if slide_def.get("url"):
                lines.append(slide_def["url"])
            set_text(slide.placeholders[1].text_frame, lines, size=18)

        else:  # prose and seeded slides
            slide = prs.slides.add_slide(layout_for(prs, "title_content"))
            slide.shapes.title.text = slide_def.get("title", "")
            body_lines = slide_def.get("bullets") or [slide_def.get("body", "")]
            body_size = 20 if len(body_lines) < 3 else 16
            if slide_def.get("image"):
                src = resolve_image(slide_def["image"], chapter_num)
                if src:
                    prepared, before, after = prepare_image(src, config)
                    stats["images"] += 1
                    stats["bytes_before"] += before
                    stats["bytes_after"] += after
                    slide.shapes.add_picture(
                        str(prepared), Inches(7.2), Inches(1.8), width=Inches(5.5)
                    )
                else:
                    stats["missing"].append(slide_def["image"])
            set_text(slide.placeholders[1].text_frame, body_lines, size=body_size)

        add_notes(slide, notes, slide_def["key"], slide_def["fingerprint"])

    return prs, stats


def render_qa_pngs(pptx_path: Path, outdir: Path):
    """Render the built deck to PNGs so it can be eyeballed without PowerPoint.
    Same soffice -> pdftoppm path Fall2026's extract_ppt.py uses."""
    soffice = shutil.which("soffice") or r"C:\Program Files\LibreOffice\program\soffice.exe"
    if not Path(soffice).exists():
        return None, "LibreOffice not found"
    outdir.mkdir(parents=True, exist_ok=True)
    for stale in outdir.glob("slide-*.png"):
        stale.unlink()
    with tempfile.TemporaryDirectory() as tmp:
        subprocess.run(
            [soffice, "--headless", "--convert-to", "pdf", "--outdir", tmp,
             str(pptx_path)],
            check=True, capture_output=True, timeout=300,
        )
        pdf = Path(tmp) / (pptx_path.stem + ".pdf")
        if not pdf.exists():
            return None, "PDF conversion produced no file"
        if shutil.which("pdftoppm"):
            subprocess.run(
                ["pdftoppm", "-png", "-r", "80", str(pdf), str(outdir / "slide")],
                check=True, capture_output=True, timeout=300,
            )
            return sorted(outdir.glob("slide-*.png")), None
        shutil.copy2(pdf, outdir / pdf.name)
        return [outdir / pdf.name], "pdftoppm not found — kept the PDF instead"


# --------------------------------------------------------------------------
# CLI
# --------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("chapter", type=int, choices=sorted(CHAPTER_SOURCES))
    parser.add_argument("--template", type=Path,
                        help="A .potx whose layouts and theme the deck inherits.")
    parser.add_argument("--outdir", type=Path, default=BUILD_DIR)
    parser.add_argument("--render", action="store_true",
                        help="Also render QA PNGs via LibreOffice.")
    args = parser.parse_args()

    source = SOURCE_DIR / CHAPTER_SOURCES[args.chapter]
    if not source.exists():
        sys.exit(f"Chapter source not found: {source}")

    lines = source.read_text(encoding="utf-8").splitlines()
    blocks = parse_blocks(lines)

    chapter_title = next(
        (b["text"] for b in blocks if b["type"] == "heading" and b["level"] == 1),
        f"Chapter {args.chapter}",
    )

    slides = build_slide_model(args.chapter, blocks, CONFIG)
    seed_path = SEED_DIR / f"ch{args.chapter:02d}.md"
    overrides = load_seed(seed_path)
    slides, applied, inserted, orphaned = apply_seed(slides, overrides)

    prs, stats = render_deck(args.chapter, chapter_title, slides, CONFIG,
                             args.template)

    args.outdir.mkdir(parents=True, exist_ok=True)
    stem = f"ch{args.chapter:02d}-{slugify(chapter_title.split(':')[-1])}"
    pptx_path = args.outdir / f"{stem}.pptx"
    prs.save(str(pptx_path))

    model = {
        "chapter": args.chapter,
        "title": chapter_title,
        "source": str(source.relative_to(ROOT)).replace("\\", "/"),
        "seed": str(seed_path.relative_to(ROOT)).replace("\\", "/")
                 if seed_path.exists() else None,
        "config": CONFIG,
        "slides": [
            {k: v for k, v in s.items() if not k.startswith("_")} for s in slides
        ],
    }
    model_path = args.outdir / f"{stem}.deck-model.json"
    model_path.write_text(json.dumps(model, indent=2, ensure_ascii=False),
                          encoding="utf-8")

    kinds = {}
    for slide in slides:
        kinds[slide["kind"]] = kinds.get(slide["kind"], 0) + 1

    print(f"Built {pptx_path.relative_to(ROOT)}")
    print(f"  {len(slides)} slides  ({pptx_path.stat().st_size / 1048576:.1f} MB)")
    print("  " + ", ".join(f"{k}={v}" for k, v in sorted(kinds.items())))
    if stats["images"]:
        print(f"  {stats['images']} images: "
              f"{stats['bytes_before'] / 1048576:.1f} MB source -> "
              f"{stats['bytes_after'] / 1048576:.1f} MB embedded "
              f"(cap {CONFIG['max_image_px']}px)")
    if stats["missing"]:
        print(f"  MISSING IMAGES ({len(stats['missing'])}):")
        for path in stats["missing"]:
            print(f"    {path}")
    if applied or inserted:
        print(f"  seed: {len(applied)} override(s), {len(inserted)} inserted slide(s)")
    if orphaned:
        print(f"  ORPHANED SEED KEYS ({len(orphaned)}) — the chapter likely moved "
              f"underneath these:")
        for key in orphaned:
            print(f"    {key}")
    print(f"  model: {model_path.relative_to(ROOT)}")

    if args.render:
        pngs, warning = render_qa_pngs(pptx_path, args.outdir / f"{stem}_qa")
        if warning:
            print(f"  QA render: {warning}")
        if pngs:
            print(f"  QA render: {len(pngs)} file(s) in "
                  f"{(args.outdir / f'{stem}_qa').relative_to(ROOT)}")


if __name__ == "__main__":
    main()
